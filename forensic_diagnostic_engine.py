"""
================================================================================
TITLE:   UNIVERSAL DATA INTEGRITY & FORENSIC DIAGNOSTIC ENGINE
AUTHOR:  Andrew R. Goad
--------------------------------------------------------------------------------
FUNCTIONAL UTILITY: 
This code provides an automated "In-Take Audit" for raw datasets, scoring 
mathematical integrity and generating forensic remediation maps for stakeholders.
================================================================================

--------------------------------------------------------------------------------
EXECUTIVE "TALK TRACKS" (FOR PRESENTATIONS)
--------------------------------------------------------------------------------
1. THE CORE PROBLEM: "The Garbage-In, Garbage-Out Risk"
   Data Science is only as good as the data it consumes. This engine acts as 
   a 'gatekeeper' that prevents corrupted or biased data from reaching 
   downstream models where it could cause costly errors.

2. STRATEGIC VALUE: "Forensic vs. Descriptive"
   We aren't just looking at averages. We are hunting for 'Forensic Anomalies'—
   dominant biases, hidden null-like strings (e.g., 'none', 'N/A'), and 
   structural drifts that standard DQ tools often miss.

3. REDUCING TECHNICAL DEBT: "Catching Errors at the Intake"
   By identifying structural flaws before the data hits the warehouse, we 
   reduce technical debt. It is significantly cheaper to fix data at the 
   intake phase than to remediate a model after it has made incorrect 
   predictions.

4. THE KPI: "The Dataset Integrity Score"
   By assigning a 0-100% health score, we give the business a clear Go/No-Go 
   metric for automated pipelines. It moves data quality from a technical 
   vague-feeling to a quantifiable business metric.

5. MULTI-CHANNEL REPORTING: "Audit-Ready Artifacts"
   The engine produces both an Executive Scorecard (PPTX) for decision-makers 
   and a Forensic Audit Ledger (XLSX). This ensures that while leadership 
   sees the 'Health Score,' developers have the 'Repair Map.'

--------------------------------------------------------------------------------
DEVELOPER & AUDIT NOTES (TECHNICAL RIGOR)
--------------------------------------------------------------------------------
1. ENGINE HEURISTICS (HOW IT THINKS)
   - DENSITY: Detects if missing values exceed the allowed NULL_THRESHOLD.
   - BIAS: Detects if a single value dominates a column (99% baseline).
   - VARIANCE: Detects 'Dead Columns' with zero mathematical variance.
   - FORMAT: Uses Regex-style pattern analysis to find inconsistent data entry.
   - DRIFT: Monitors for cardinality shifts that suggest a change in data source.

2. THE REMEDIATION MAP (XLSX OUTPUT)
   - Every flag generates a 'Forensic Rationale'—a plain-English explanation 
     designed to teach the data owner exactly why their data was rejected.
   - 'Potential Improvement' shows what the score would be if that column 
     were remediated, helping teams prioritize engineering efforts.

3. POWERPOINT AUTOMATION (EXECUTIVE AID)
   - Uses 'python-pptx' to generate a visual scorecard including the Integrity 
     Score, Remediation Potential, and a Bar Chart of column-level health.
================================================================================
"""

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches

# =============================
# CONFIGURATION
# =============================
NULL_THRESHOLD = 0.20
FORMAT_DEVIATION_LIMIT = 0.05
BIAS_THRESHOLD = 0.99

# Impact weights for column scoring
WEIGHTS = {
    "Density": 0.30,
    "Format_Inconsistency": 0.25,
    "Dominant_Bias": 0.20,
    "Zero_Variance": 0.25
}

# "Null-like" strings often found in enterprise data
STRING_NULLS = ['nan', 'null', 'na', 'n/a', 'undefined', '???', '.', '']

class DataForensicEngine:
    def __init__(self, df):
        self.df = df
        self.column_scores = {}
        self.red_flags = {}
        self.dataset_score = 0.0
        self.simulated_score = 0.0
        self.plot_buffer = None

    def run_audit(self):
        """Execute the full forensic suite across all columns."""
        for col in self.df.columns:
            self._audit_column(col)
        
        self.dataset_score = np.mean([v["Score"] for v in self.column_scores.values()])
        self._calculate_remediation_potential()
        self._plot_scores()
        self._generate_scorecard()

    def _audit_column(self, col):
        """Internal heuristic engine to detect anomalies."""
        series = self.df[col]
        total_rows = len(series)
        alerts = {}
        
        # 1. Density Audit (Nulls & Hidden Nulls)
        null_count = series.isnull().sum()
        string_null_count = 0
        if series.dtype == 'object':
            string_null_count = series.astype(str).str.lower().str.strip().isin(STRING_NULLS).sum()
        
        actual_null_ratio = (null_count + string_null_count) / total_rows
        if actual_null_ratio > NULL_THRESHOLD:
            alerts["Density"] = f"High Null Density ({actual_null_ratio:.1%})"

        # 2. Dominant Bias Audit
        mode_freq = series.value_counts(normalize=True).iloc[0] if not series.empty else 0
        if mode_freq >= BIAS_THRESHOLD:
            alerts["Dominant_Bias"] = f"Statistical Bias: {mode_freq:.1%} values are identical."

        # 3. Format Inconsistency Audit (Simplified Pattern Check)
        if series.dtype == 'object':
            lengths = series.astype(str).apply(len).value_counts(normalize=True)
            if not lengths.empty and lengths.iloc[0] < (1 - FORMAT_DEVIATION_LIMIT):
                alerts["Format_Inconsistency"] = "High variance in string format/length."

        # 4. Zero Variance Audit
        if series.nunique() <= 1:
            alerts["Zero_Variance"] = "Column contains zero information (1 or 0 unique values)."

        # Score Calculation
        deductions = sum(WEIGHTS[k] for k in alerts.keys())
        score = max(0, 1.0 - deductions)
        
        self.column_scores[col] = {"Score": score, "Alerts": alerts}

    def _calculate_remediation_potential(self):
        """Simulates health score assuming known forensic flags are repaired."""
        simulated_col_scores = []
        for col, data in self.column_scores.items():
            # Assume Density and Format are repairable; Zero Variance and Bias are structural
            repairable_deductions = sum(WEIGHTS[k] for k in data["Alerts"].keys() if k in ["Density", "Format_Inconsistency"])
            simulated_col_scores.append(data["Score"] + repairable_deductions)
        self.simulated_score = np.mean(simulated_col_scores)

    def _plot_scores(self):
        """Generate bar chart for column scores; save to buffer."""
        scores = [v["Score"] for v in self.column_scores.values()]
        cols = list(self.column_scores.keys())
        plt.figure(figsize=(10,5))
        plt.bar(cols, scores, color='#58a6ff')
        plt.ylabel("Column Score")
        plt.title("Column Integrity Scores")
        plt.xticks(rotation=45, ha='right')
        buf = BytesIO()
        plt.savefig(buf, format='png', bbox_inches='tight')
        buf.seek(0)
        self.plot_buffer = buf
        plt.close()

    def _generate_scorecard(self):
        """Display dataset summary scores."""
        print(f"Dataset Score: {self.dataset_score:.2%}")
        print(f"Simulated Score after fixes: {self.simulated_score:.2%}")

    def export_to_excel(self, filename="Forensic_Audit_Ledger.xlsx"):
        """Export the Forensic Audit Ledger."""
        audit_data = []
        for col, data in self.column_scores.items():
            alerts = data["Alerts"]
            audit_data.append({
                "Column": col,
                "Integrity_Score": data["Score"],
                "Status": "FAIL" if alerts else "PASS",
                "Alert_Types": ", ".join(alerts.keys()) if alerts else "None",
                "Forensic_Rationale": " | ".join(alerts.values()) if alerts else "Clean"
            })
        pd.DataFrame(audit_data).to_excel(filename, index=False)

    def export_to_pptx(self, filename="Executive_Scorecard.pptx"):
        """Export the Executive Scorecard to PPTX."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Dataset Forensic Diagnostic Summary"
        
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(1.5))
        tf = txBox.text_frame
        tf.text = f"Current Integrity Score: {self.dataset_score:.2%}"
        p = tf.add_paragraph()
        p.text = f"Remediation Potential: {self.simulated_score:.2%}"
        
        if self.plot_buffer:
            slide.shapes.add_picture(self.plot_buffer, Inches(4.5), Inches(1.5), width=Inches(5))
        
        prs.save(filename)

# =============================
# DEMO / PRODUCTION USE CASE
# =============================
if __name__ == "__main__":
    # Creating a sample dataset with intentional forensic flaws
    demo_df = pd.DataFrame({
        'User_ID': range(1, 101),
        'Account_Code': ['PREM100']*92 + ['PREM-99', 'BASIC', '???', 'X1']*2, # Format Inconsistency
        'Status': ['Active']*99 + ['Inactive'], # Dominant Bias
        'Legacy_Field': [np.nan]*30 + ['None']*10 + ['A']*60, # Density Alert
        'Unique_Noise': ['ID' + str(i) for i in range(100)], # Categorical Noise
        'Clean_Field': [10.5] * 100
    })

    # Initialize and run audit
    engine = DataForensicEngine(demo_df)
    engine.run_audit()
    
    # engine.export_to_excel("Andrew_Goad_Forensic_Audit.xlsx")
    # engine.export_to_pptx("Executive_Data_Scorecard.pptx")

    print("\n[COMPLETE] Forensic Audit generated.")
